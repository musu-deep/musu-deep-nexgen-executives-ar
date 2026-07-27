from __future__ import annotations

import csv
import io
import json
import re
import statistics
from collections import Counter
from pathlib import Path
from typing import Any

MAX_TEXT_CHARS = 180_000
MAX_FILE_BYTES = 4_000_000

ARABIC_DIGITS = str.maketrans("٠١٢٣٤٥٦٧٨٩۰۱۲۳۴۵۶۷۸۹", "01234567890123456789")

CURRENCY_MARKERS = {
    "SAR": ("sar", "ر.س", "ريال", "ريال سعودي"),
    "USD": ("usd", "$", "دولار"),
    "AED": ("aed", "درهم"),
    "EUR": ("eur", "€", "يورو"),
    "GBP": ("gbp", "£", "جنيه إسترليني"),
    "XOF": ("xof", "fcfa", "cfa"),
}

COST_WORDS = (
    "تكلفة", "تكاليف", "cost", "expense", "مصروف", "مواد", "عمالة",
    "توريد", "تنفيذ", "مصاريف مباشرة", "direct cost",
)
QUOTE_WORDS = (
    "السعر", "سعر", "عرض", "قيمة العرض", "quotation", "quote", "bid",
    "proposal", "price", "commercial offer",
)
TOTAL_WORDS = (
    "الإجمالي", "اجمالي", "المجموع", "القيمة الإجمالية", "total", "grand total",
    "subtotal", "صافي القيمة",
)
RISK_WORDS = (
    "مخاطر", "تأخير", "تعثر", "غرامة", "مطالبة", "نزاع", "عاجل", "حرج",
    "risk", "delay", "penalty", "claim", "dispute", "breach", "critical",
)
OBLIGATION_WORDS = (
    "يلتزم", "يجب", "يتعين", "اعتماد", "تسليم", "موعد", "دفعة", "مرحلة",
    "shall", "must", "deliver", "deadline", "approve", "payment", "milestone",
)


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "cp1256", "windows-1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _clean_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()[:MAX_TEXT_CHARS]


def _extract_pdf(data: bytes) -> tuple[str, dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    text = "\n".join((page.extract_text() or "") for page in reader.pages[:120])
    return text, {"pages": len(reader.pages), "parser": "pypdf"}


def _extract_docx(data: bytes) -> tuple[str, dict[str, Any]]:
    from docx import Document

    document = Document(io.BytesIO(data))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                lines.append("\t".join(values))
    return "\n".join(lines), {"paragraphs": len(document.paragraphs), "tables": len(document.tables), "parser": "python-docx"}


def _extract_xlsx(data: bytes) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    row_count = 0
    for worksheet in workbook.worksheets[:20]:
        lines.append(f"[ورقة: {worksheet.title}]")
        for row in worksheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                lines.append("\t".join(values))
                row_count += 1
            if row_count >= 8000:
                break
        if row_count >= 8000:
            break
    return "\n".join(lines), {"sheets": len(workbook.worksheets), "rows": row_count, "parser": "openpyxl"}


def _extract_pptx(data: bytes) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[شريحة {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                lines.append(str(shape.text).strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        lines.append("\t".join(values))
    return "\n".join(lines), {"slides": len(presentation.slides), "parser": "python-pptx"}


def extract_file_text(filename: str, content_type: str, data: bytes) -> dict[str, Any]:
    if not data:
        raise ValueError("الملف فارغ")
    if len(data) > MAX_FILE_BYTES:
        raise ValueError("حجم الملف يتجاوز 4 ميجابايت. قسّم الملف أو ارفع نسخة مضغوطة نصيًا.")

    suffix = Path(filename or "file").suffix.lower()
    metadata: dict[str, Any] = {"parser": "text", "extension": suffix, "size_bytes": len(data)}
    warning = ""
    try:
        if suffix == ".pdf" or content_type == "application/pdf":
            text, extra = _extract_pdf(data)
        elif suffix == ".docx":
            text, extra = _extract_docx(data)
        elif suffix in {".xlsx", ".xlsm"}:
            text, extra = _extract_xlsx(data)
        elif suffix == ".pptx":
            text, extra = _extract_pptx(data)
        elif suffix == ".json":
            payload = json.loads(_decode_text(data))
            text = json.dumps(payload, ensure_ascii=False, indent=2)
            extra = {"parser": "json"}
        elif suffix == ".csv":
            decoded = _decode_text(data)
            rows = list(csv.reader(io.StringIO(decoded)))
            text = "\n".join("\t".join(row) for row in rows[:10000])
            extra = {"rows": len(rows), "parser": "csv"}
        elif suffix in {".txt", ".md", ".html", ".htm", ".xml", ".tsv"} or content_type.startswith("text/"):
            text = _decode_text(data)
            extra = {"parser": "text"}
        else:
            text = _decode_text(data)
            extra = {"parser": "fallback-text"}
            warning = "نوع الملف غير معروف بالكامل؛ تمت محاولة قراءة محتواه كنص."
    except Exception as exc:
        text = _decode_text(data)
        extra = {"parser": "fallback-text"}
        warning = f"تعذر استخدام القارئ المتخصص؛ استُخدمت القراءة النصية الاحتياطية: {str(exc)[:160]}"

    text = _clean_text(text)
    if not text:
        warning = warning or "لم يُستخرج نص قابل للتحليل. قد يكون الملف صورة ممسوحة ضوئيًا ويحتاج OCR."
    metadata.update(extra)
    metadata["characters"] = len(text)
    metadata["warning"] = warning
    return {"text": text, "metadata": metadata}


def _normalise_number(raw: str) -> float | None:
    value = raw.translate(ARABIC_DIGITS).replace(" ", "").replace(",", "")
    try:
        return float(value)
    except ValueError:
        return None


def _detect_currency(context: str) -> str | None:
    lowered = context.lower()
    for code, markers in CURRENCY_MARKERS.items():
        if any(marker in lowered for marker in markers):
            return code
    return None


def _amount_kind(context: str) -> str:
    lowered = context.lower()
    if any(word in lowered for word in TOTAL_WORDS):
        return "total"
    if any(word in lowered for word in COST_WORDS):
        return "cost"
    if any(word in lowered for word in QUOTE_WORDS):
        return "quote"
    return "amount"


def extract_money_amounts(text: str) -> list[dict[str, Any]]:
    normalised = text.translate(ARABIC_DIGITS)
    pattern = re.compile(r"(?<!\d)(\d{1,3}(?:[ ,]\d{3})+(?:\.\d+)?|\d{3,}(?:\.\d+)?)(?!\d)")
    results: list[dict[str, Any]] = []
    seen: set[tuple[float, str, str | None]] = set()
    for match in pattern.finditer(normalised):
        value = _normalise_number(match.group(1))
        if value is None or value <= 0:
            continue
        if 1900 <= value <= 2100 and match.group(1).isdigit():
            continue
        context = normalised[max(0, match.start() - 55): min(len(normalised), match.end() + 55)]
        currency = _detect_currency(context)
        kind = _amount_kind(context)
        if value < 1000 and not currency and kind == "amount":
            continue
        key = (round(value, 2), kind, currency)
        if key in seen:
            continue
        seen.add(key)
        results.append({
            "value": round(value, 2),
            "currency": currency,
            "kind": kind,
            "context": re.sub(r"\s+", " ", context).strip()[:150],
        })
    return results[:250]


def _extract_dates(text: str) -> list[str]:
    patterns = (
        r"\b\d{4}-\d{1,2}-\d{1,2}\b",
        r"\b\d{1,2}/\d{1,2}/\d{2,4}\b",
        r"\b\d{1,2}-\d{1,2}-\d{2,4}\b",
    )
    found: list[str] = []
    for pattern in patterns:
        found.extend(re.findall(pattern, text.translate(ARABIC_DIGITS)))
    return list(dict.fromkeys(found))[:12]


def analyse_document_text(title: str, description: str, text: str, purpose: str = "general") -> dict[str, Any]:
    combined = _clean_text("\n".join(value for value in (title, description, text) if value))
    lowered = combined.lower()
    lines = [line.strip() for line in combined.splitlines() if len(line.strip()) >= 12]
    risks = [word for word in RISK_WORDS if word in lowered]
    obligation_lines = [line for line in lines if any(word in line.lower() for word in OBLIGATION_WORDS)][:8]
    amounts = extract_money_amounts(combined)
    risk_level = "high" if len(risks) >= 3 else "medium" if risks else "low"
    excerpt = description.strip() or " ".join(lines[:5])
    if not excerpt:
        excerpt = "تم استلام الملف، لكن لم يُستخرج منه نص كافٍ لبناء ملخص تفصيلي."

    purpose_labels = {
        "advisory": "ملف استشاري",
        "opportunity": "ملف فرصة أو عرض",
        "pricing": "ملف تسعير أو عرض مالي",
        "general": "مستند مؤسسي",
    }
    summary = f"تمت قراءة {purpose_labels.get(purpose, 'مستند مؤسسي')} بعنوان «{title}». {excerpt[:650]}"
    return {
        "summary": summary,
        "parties": ["مجموعة اراك للتنمية", "الجهة أو العميل الوارد في الملف"],
        "dates": _extract_dates(combined),
        "obligations": obligation_lines or [
            "مراجعة نطاق الملف وتحديد المسؤول التنفيذي.",
            "التحقق من المواعيد والقيم والافتراضات قبل اعتماد القرار.",
        ],
        "risks": ([{"level": risk_level, "risk": f"رُصدت مؤشرات تحتاج مراجعة: {،.join(risks[:6])}."}] if risks else []),
        "important_clauses": [
            "النطاق والمخرجات",
            "القيمة أو التكلفة وشروط الدفع",
            "المواعيد والمسؤوليات",
            "الافتراضات والاستثناءات",
        ],
        "monetary_values": amounts[:30],
        "risk_level": risk_level,
        "suggested_task": {
            "title": f"مراجعة وتحليل: {title}",
            "description": "التحقق من المخرجات المستخرجة وربطها بصاحب القرار والإجراء التالي.",
            "priority": "high" if risk_level == "high" else "medium",
            "sector": "corporate",
        },
        "suggested_meeting": {
            "title": f"مراجعة تنفيذية للملف: {title}",
            "duration_minutes": 30,
            "reason": "مراجعة القيم والمخاطر والمسؤوليات والقرار المطلوب.",
        },
        "generated_by": "ARAAK File Intelligence Engine",
    }


def _median(values: list[float]) -> float:
    return round(float(statistics.median(values)), 2) if values else 0.0


def analyse_pricing_documents(
    documents: list[dict[str, Any]],
    target_margin: float = 18,
    overhead_rate: float = 8,
    risk_rate: float = 6,
    win_strength: float = 75,
) -> dict[str, Any]:
    file_summaries: list[dict[str, Any]] = []
    file_totals: list[float] = []
    file_costs: list[float] = []
    currency_votes: list[str] = []

    for item in documents:
        amounts = extract_money_amounts(item.get("text", ""))
        total_values = [entry["value"] for entry in amounts if entry["kind"] in {"total", "quote"}]
        cost_values = [entry["value"] for entry in amounts if entry["kind"] == "cost"]
        all_values = [entry["value"] for entry in amounts]
        representative = max(total_values) if total_values else max(all_values) if all_values else 0
        representative_cost = max(cost_values) if cost_values else 0
        if representative:
            file_totals.append(representative)
        if representative_cost:
            file_costs.append(representative_cost)
        currency_votes.extend(entry["currency"] for entry in amounts if entry.get("currency"))
        file_summaries.append({
            "document_id": item.get("document_id"),
            "filename": item.get("filename"),
            "title": item.get("title"),
            "amount_count": len(amounts),
            "representative_price": round(representative, 2),
            "detected_cost": round(representative_cost, 2),
            "minimum": round(min(all_values), 2) if all_values else 0,
            "maximum": round(max(all_values), 2) if all_values else 0,
            "amounts": amounts[:25],
            "warning": item.get("metadata", {}).get("warning", ""),
        })

    benchmark = _median(file_totals)
    detected_direct_cost = _median(file_costs)
    currency = Counter(currency_votes).most_common(1)[0][0] if currency_votes else "SAR"

    if not benchmark and detected_direct_cost:
        benchmark = detected_direct_cost / max(0.25, 1 - (target_margin + overhead_rate + risk_rate) / 100)

    estimated_cost = detected_direct_cost or (benchmark * 0.72 if benchmark else 0)
    overhead = estimated_cost * max(0, overhead_rate) / 100
    risk = (estimated_cost + overhead) * max(0, risk_rate) / 100
    full_cost = estimated_cost + overhead + risk
    floor_price = full_cost / 0.90 if full_cost else 0
    cost_based = full_cost / max(0.05, 1 - max(0, min(target_margin, 70)) / 100) if full_cost else 0

    if benchmark:
        balanced_seed = cost_based * 0.65 + benchmark * 0.35 if cost_based else benchmark * 0.98
        balanced = max(floor_price, min(max(balanced_seed, benchmark * 0.94), benchmark * 1.06))
        aggressive = max(floor_price, min(balanced * 0.97, benchmark * 0.965))
        premium = max(balanced * 1.055, benchmark * 1.035)
    else:
        balanced = cost_based
        aggressive = max(floor_price, balanced * 0.97)
        premium = balanced * 1.055

    gap = benchmark - balanced if benchmark else 0
    price_advantage = (gap / benchmark * 100) if benchmark else 0
    win_probability = max(5, min(98, win_strength + price_advantage * 1.8 - risk_rate * 0.55))
    confidence = 35 + min(35, len(file_summaries) * 9) + (15 if benchmark else 0) + (15 if detected_direct_cost else 0)
    confidence = min(95, confidence)

    recommendation = (
        "اعتماد العرض المتوازن مع توضيح القيمة غير السعرية، وإبقاء الباقة الهجومية للتفاوض النهائي فقط."
        if benchmark and balanced <= benchmark
        else "إعادة هندسة النطاق أو التكلفة قبل التقديم؛ السعر المحسوب أعلى من المؤشر المرجعي أو لا توجد بيانات سوقية كافية."
    )
    if not file_totals:
        recommendation = "لم تُستخرج أسعار نهائية موثوقة. راجع الملفات أو أدخل الأسعار المرجعية يدويًا ثم أعد التحليل."

    return {
        "currency": currency,
        "confidence": round(confidence),
        "benchmark": {
            "minimum": round(min(file_totals), 2) if file_totals else 0,
            "maximum": round(max(file_totals), 2) if file_totals else 0,
            "median": round(benchmark, 2),
            "sample_size": len(file_totals),
        },
        "cost": {
            "detected_direct_cost": round(detected_direct_cost, 2),
            "estimated_direct_cost": round(estimated_cost, 2),
            "overhead": round(overhead, 2),
            "risk_reserve": round(risk, 2),
            "full_cost": round(full_cost, 2),
            "safe_floor": round(floor_price, 2),
        },
        "offer_options": [
            {"key": "aggressive", "name": "عرض هجومي", "price": round(aggressive, 2), "purpose": "رفع احتمالية الفوز مع أقل مساحة تفاوض"},
            {"key": "balanced", "name": "عرض متوازن", "price": round(balanced, 2), "purpose": "أفضل توازن بين التنافسية والهامش"},
            {"key": "premium", "name": "عرض قيمة مضافة", "price": round(premium, 2), "purpose": "نطاق أوسع وقيمة استراتيجية وضمانات أعلى"},
        ],
        "recommended_price": round(balanced, 2),
        "negotiation_floor": round(max(floor_price, aggressive), 2),
        "competitive_gap": round(gap, 2),
        "price_advantage_percent": round(price_advantage, 1),
        "win_probability": round(win_probability),
        "recommendation": recommendation,
        "files": file_summaries,
        "model_patch": {
            "competitorPrice": round(benchmark, 2) if benchmark else 0,
            "directCost": round(detected_direct_cost, 2) if detected_direct_cost else 0,
            "winStrength": round(win_probability),
        },
    }
